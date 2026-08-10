from io import BytesIO
from pathlib import Path

from pptx import Presentation

from ..models import NormalizedUnit, ParsedDocument
from .base import DocumentParser


class PptxParser(DocumentParser):
    def parse(
        self,
        data: bytes,
        filename: str,
        mime_type: str | None,
    ) -> ParsedDocument:
        presentation = Presentation(
            BytesIO(data),
        )

        units: list[NormalizedUnit] = []

        for slide_index, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            texts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text.strip()

                if text:
                    texts.append(text)

            slide_text = "\n".join(texts)

            units.append(
                NormalizedUnit(
                    index=slide_index,
                    kind="slide",
                    label=f"Slide {slide_index}",
                    text=slide_text,
                    metadata={
                        "slide_number": slide_index,
                    },
                )
            )

        full_text = "\n\n".join(
            unit.text
            for unit in units
            if unit.text
        )

        return ParsedDocument(
            filename=filename,
            extension=Path(filename).suffix.lower(),
            mime_type=mime_type,
            parser="python-pptx",
            units=units,
            full_text=full_text,
            metadata={
                "slide_count": len(units),
            },
        )